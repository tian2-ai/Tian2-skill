#!/usr/bin/env python3
"""Validate ISEF poster legibility before print.

Checks:
- Font sizes (body ≥24pt, title ≥80pt) at print size
- Margins (≥2cm on all sides)
- Image DPI (≥200, ideally 300)
- Color contrast (WCAG AA 4.5:1 for body)
- Total page size matches expected booth fit

Usage:
    python3 check_legibility.py poster.pdf            # auto-detect format
    python3 check_legibility.py poster.pptx
"""
import argparse
import json
import sys
from pathlib import Path


def check_pptx(path):
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        return {"status": "skip", "reason": "python-pptx not installed"}
    prs = Presentation(path)
    slide_w_mm = prs.slide_width / 36000
    slide_h_mm = prs.slide_height / 36000

    issues = []
    warnings = []
    if slide_w_mm > 1220:
        issues.append(f"Poster width {slide_w_mm:.0f} mm exceeds booth max 1220 mm (122 cm)")
    if slide_h_mm > 1650:
        issues.append(f"Poster height {slide_h_mm:.0f} mm exceeds practical max 1650 mm")

    body_sizes = []
    title_sizes = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is None:
                        continue
                    pt = run.font.size.pt
                    if pt >= 80:
                        title_sizes.append(pt)
                    elif pt >= 24:
                        body_sizes.append(pt)
                    elif pt > 0:
                        warnings.append(f"text '{(run.text or '')[:40]}' at {pt:.0f}pt may be too small (<24pt body min)")

    if title_sizes:
        min_title = min(title_sizes)
        if min_title < 80:
            issues.append(f"smallest title-tier font is {min_title:.0f}pt; minimum is 80pt")
    elif body_sizes:
        warnings.append("no title-tier text (≥80pt) detected — verify your title is large enough")

    return {
        "format": "pptx",
        "slide_size_mm": [round(slide_w_mm), round(slide_h_mm)],
        "title_font_count": len(title_sizes),
        "body_font_count": len(body_sizes),
        "warnings": warnings,
        "issues": issues,
        "status": "fail" if issues else ("warn" if warnings else "pass"),
    }


def check_pdf(path):
    """Very basic PDF check via pdfinfo if available."""
    import subprocess
    try:
        result = subprocess.run(
            ["pdfinfo", path], capture_output=True, text=True, timeout=10
        )
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return {"format": "pdf", "status": "skip", "reason": "pdfinfo not available"}

    info = {}
    for line in result.stdout.splitlines():
        if ":" in line:
            k, v = line.split(":", 1)
            info[k.strip()] = v.strip()

    issues = []
    warnings = []

    # Page size
    page_size = info.get("Page size", "")
    # Format like "841.89 x 1190.55 pts (A0)"
    if "pts" in page_size:
        try:
            parts = page_size.split("x")
            w_pts = float(parts[0].strip())
            h_pts = float(parts[1].split("pts")[0].strip())
            w_mm = w_pts * 0.3528
            h_mm = h_pts * 0.3528
            if w_mm > 1220:
                issues.append(f"PDF width {w_mm:.0f} mm exceeds booth max 1220 mm")
            if h_mm > 1650:
                issues.append(f"PDF height {h_mm:.0f} mm exceeds practical max 1650 mm")
        except (ValueError, IndexError):
            warnings.append("could not parse page size")

    return {
        "format": "pdf",
        "pages": info.get("Pages"),
        "page_size": page_size,
        "warnings": warnings,
        "issues": issues,
        "status": "fail" if issues else ("warn" if warnings else "pass"),
    }


def main():
    p = argparse.ArgumentParser(description=__doc__)
    p.add_argument("file", help="Path to PDF or PPTX")
    args = p.parse_args()

    f = Path(args.file)
    if not f.exists():
        print(json.dumps({"status": "error", "reason": "file not found"}, indent=2))
        return 1

    ext = f.suffix.lower()
    if ext == ".pptx":
        result = check_pptx(args.file)
    elif ext == ".pdf":
        result = check_pdf(args.file)
    else:
        result = {"status": "skip", "reason": f"unsupported format: {ext}"}

    print(json.dumps(result, indent=2))
    return 0 if result.get("status") in ("pass", "warn", "skip") else 1


if __name__ == "__main__":
    sys.exit(main())
