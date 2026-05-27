#!/usr/bin/env python3
"""Build a print-ready ISEF poster PPTX from a content YAML.

Requires:
    python3 -m pip install python-pptx pyyaml

Usage:
    python3 build_pptx.py --content default-a0-content.yaml --size A0 --out my-poster.pptx

Sizes:
    A0          → 841 × 1189 mm (portrait)
    A0-landscape → 1189 × 841 mm
    48x36       → 48 × 36 in (landscape; US convention)
    36x48       → 36 × 48 in (portrait; US convention)
    custom W H  → W cm × H cm (must fit booth 122×165)
"""
import argparse
import json
import sys
from pathlib import Path


SIZES_MM = {
    "A0": (841, 1189),
    "A0-landscape": (1189, 841),
    "48x36": (1219, 914),    # 48"×36" = 121.9×91.4 cm
    "36x48": (914, 1219),
    "44x44": (1118, 1118),
}


def load_content(path):
    try:
        import yaml
    except ImportError:
        print("error: pyyaml is required. Install with: python3 -m pip install pyyaml", file=sys.stderr)
        sys.exit(2)
    with open(path) as f:
        return yaml.safe_load(f)


def validate_size(size_str):
    if size_str in SIZES_MM:
        return SIZES_MM[size_str]
    # parse "WxH" custom
    if "x" in size_str.lower():
        parts = size_str.lower().split("x")
        try:
            w, h = float(parts[0]), float(parts[1])
        except (ValueError, IndexError):
            raise ValueError(f"Unparseable size: {size_str}")
        # Treat as cm; convert to mm
        return int(w * 10), int(h * 10)
    raise ValueError(f"Unknown size: {size_str}")


def build_pptx(content, size_mm, out_path):
    try:
        from pptx import Presentation
        from pptx.util import Mm, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        print("error: python-pptx is required. Install with: python3 -m pip install python-pptx", file=sys.stderr)
        sys.exit(2)

    width_mm, height_mm = size_mm
    if width_mm > 1220 or height_mm > 1650:
        print(f"warning: size {width_mm}×{height_mm} mm exceeds ISEF booth max (1220×1650 mm)", file=sys.stderr)

    prs = Presentation()
    prs.slide_width = Mm(width_mm)
    prs.slide_height = Mm(height_mm)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    theme = content.get("theme", {})
    primary = _hex_to_rgb(theme.get("primary_color", "#0B3D91"))
    accent = _hex_to_rgb(theme.get("accent_color", "#FFA500"))

    proj = content["project"]

    margin = 30  # mm
    inner_w = width_mm - 2 * margin

    # Banner
    banner_h = 110
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Mm(0), Mm(0), Mm(width_mm), Mm(banner_h))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(*primary)
    banner.line.fill.background()

    # Title
    tb = slide.shapes.add_textbox(Mm(margin), Mm(10), Mm(inner_w), Mm(banner_h - 20))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = proj.get("title", "")
    r.font.size = Pt(100)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = f"{proj.get('authors', '')} • {proj.get('affiliation', '')} • {proj.get('category', '')} • {proj.get('project_id', '')}"
    r2.font.size = Pt(36)
    r2.font.color.rgb = RGBColor(255, 255, 255)

    # Three-column body region
    y_cursor = banner_h + margin
    col_w = (inner_w - 2 * 20) / 3
    col_x = [margin, margin + col_w + 20, margin + 2 * (col_w + 20)]

    # Abstract — full width
    abstract = content.get("abstract", "")
    if abstract:
        ab_h = 180
        ab = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Mm(margin), Mm(y_cursor), Mm(inner_w), Mm(ab_h))
        ab.fill.solid()
        ab.fill.fore_color.rgb = RGBColor(*accent)
        ab.line.fill.background()
        ab_tb = slide.shapes.add_textbox(Mm(margin + 10), Mm(y_cursor + 10), Mm(inner_w - 20), Mm(ab_h - 20))
        ab_tf = ab_tb.text_frame
        ab_tf.word_wrap = True
        p = ab_tf.paragraphs[0]
        r = p.add_run()
        r.text = "ABSTRACT"
        r.font.size = Pt(36)
        r.font.bold = True
        r.font.color.rgb = RGBColor(255, 255, 255)
        p2 = ab_tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = abstract.strip()
        r2.font.size = Pt(24)
        r2.font.color.rgb = RGBColor(255, 255, 255)
        y_cursor += ab_h + 20

    # Sections — flow across columns
    sections = content.get("sections", [])
    col_idx = 0
    col_y = [y_cursor, y_cursor, y_cursor]
    for sec in sections:
        x = col_x[col_idx]
        y = col_y[col_idx]
        body_text = sec.get("body", "")
        figure = sec.get("figure")
        caption = sec.get("caption", "")
        heading = sec.get("heading", "")
        section_h = 0
        if heading:
            h_tb = slide.shapes.add_textbox(Mm(x), Mm(y), Mm(col_w), Mm(20))
            h_tf = h_tb.text_frame
            r = h_tf.paragraphs[0].add_run()
            r.text = heading
            r.font.size = Pt(48)
            r.font.bold = True
            r.font.color.rgb = RGBColor(*primary)
            section_h += 25
        if body_text:
            b_tb = slide.shapes.add_textbox(Mm(x), Mm(y + section_h), Mm(col_w), Mm(120))
            b_tf = b_tb.text_frame
            b_tf.word_wrap = True
            r = b_tf.paragraphs[0].add_run()
            r.text = body_text.strip()
            r.font.size = Pt(28)
            section_h += 130
        if figure and Path(figure).exists():
            slide.shapes.add_picture(figure, Mm(x), Mm(y + section_h), width=Mm(col_w))
            section_h += int(col_w * 0.6)  # rough aspect estimate
            if caption:
                c_tb = slide.shapes.add_textbox(Mm(x), Mm(y + section_h), Mm(col_w), Mm(20))
                c_tf = c_tb.text_frame
                c_tf.word_wrap = True
                r = c_tf.paragraphs[0].add_run()
                r.text = caption
                r.font.size = Pt(20)
                r.font.italic = True
                section_h += 22
        col_y[col_idx] = y + section_h + 20
        # Move to next column if this one is full
        if col_y[col_idx] > height_mm - 200:
            col_idx = (col_idx + 1) % 3

    # Footer — references + ack + AI disclosure
    footer_y = height_mm - 180
    refs = content.get("references", [])
    if refs:
        ref_tb = slide.shapes.add_textbox(Mm(margin), Mm(footer_y), Mm(inner_w * 0.5), Mm(160))
        ref_tf = ref_tb.text_frame
        ref_tf.word_wrap = True
        r = ref_tf.paragraphs[0].add_run()
        r.text = "REFERENCES"
        r.font.size = Pt(24)
        r.font.bold = True
        r.font.color.rgb = RGBColor(*primary)
        for ref in refs:
            p = ref_tf.add_paragraph()
            rr = p.add_run()
            rr.text = "• " + ref
            rr.font.size = Pt(16)

    acks = content.get("acknowledgments", "")
    aiuse = content.get("ai_use_disclosure", "")
    if acks or aiuse:
        ack_tb = slide.shapes.add_textbox(Mm(margin + inner_w * 0.55), Mm(footer_y), Mm(inner_w * 0.45), Mm(160))
        ack_tf = ack_tb.text_frame
        ack_tf.word_wrap = True
        if acks:
            r = ack_tf.paragraphs[0].add_run()
            r.text = "ACKNOWLEDGMENTS"
            r.font.size = Pt(20)
            r.font.bold = True
            r.font.color.rgb = RGBColor(*primary)
            p = ack_tf.add_paragraph()
            rr = p.add_run()
            rr.text = acks.strip()
            rr.font.size = Pt(16)
        if aiuse:
            p = ack_tf.add_paragraph()
            r = p.add_run()
            r.text = "\nAI-USE DISCLOSURE"
            r.font.size = Pt(20)
            r.font.bold = True
            r.font.color.rgb = RGBColor(*primary)
            p2 = ack_tf.add_paragraph()
            rr = p2.add_run()
            rr.text = aiuse.strip()
            rr.font.size = Pt(14)
            rr.font.italic = True

    prs.save(out_path)
    print(f"wrote {out_path}", file=sys.stderr)
    print(json.dumps({"out": str(out_path), "size_mm": list(size_mm), "sections": len(sections)}, indent=2))


def _hex_to_rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def main():
    p = argparse.ArgumentParser(description=__doc__)
    p.add_argument("--content", required=True, help="Path to content YAML")
    p.add_argument("--size", default="A0")
    p.add_argument("--out", required=True)
    args = p.parse_args()

    content = load_content(args.content)
    size_mm = validate_size(args.size)
    build_pptx(content, size_mm, args.out)


if __name__ == "__main__":
    main()
